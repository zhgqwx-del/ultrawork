#!/usr/bin/env python3
"""deckcraft negative-sample selftest — proves the gates CATCH violations.

The bundled example (examples/ai-coding-pilot) is the positive control: it can
only prove compliant input passes. This script synthesizes violating projects
in a temp dir and asserts each gate rejects them with the expected error code.
Run after any change to skills/builtin/deckcraft/scripts/ or templates:

    python3 scripts/deckcraft-selftest.py            # full run (needs Chrome for probe cases)
    python3 scripts/deckcraft-selftest.py --no-chrome  # skip browser-dependent cases

Exit 0 = all cases behave, 1 = a gate failed to catch (or a positive case failed).
"""
from __future__ import annotations

import json
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

REPO = Path(__file__).resolve().parent.parent
SCRIPTS = REPO / "skills" / "builtin" / "deckcraft" / "scripts"
PY = sys.executable  # never hardcode python3 vs python (Windows)

TOKENS = (":root{--c-bg:#ffffff;--c-fg:#1a1a1a;--c-primary:#1a1a2e;--c-accent:#c8501e;"
          "--c-muted:#6a6a6a;--c-on-dark:#f5f5f5;--c-line:#dddddd;"
          "--fs-hero:64px;--fs-h1:40px;--fs-h3:22px;--fs-body:17px;--fs-caption:12px}")

PASS = 0
FAIL = 0


def run(script: str, *args: str) -> tuple[int, str]:
    r = subprocess.run([PY, str(SCRIPTS / script), *args],
                       capture_output=True, text=True, encoding="utf-8",
                       errors="replace", timeout=180)
    return r.returncode, (r.stdout or "") + (r.stderr or "")


def project(td: str, pages: dict[int, str], outline: dict | None = None) -> Path:
    proj = Path(td) / "proj"
    (proj / "pages").mkdir(parents=True)
    (proj / "tokens.css").write_text(TOKENS, encoding="utf-8")
    if outline is not None:
        (proj / "outline.json").write_text(
            json.dumps(outline, ensure_ascii=False), encoding="utf-8")
    for n, frag in pages.items():
        (proj / "pages" / f"page-{n:02d}.html").write_text(frag, encoding="utf-8")
    return proj


def page(inner: str, layout: str = "S03", rhythm: str = "dense") -> str:
    return (f'<section class="slide" data-layout="{layout}" data-rhythm="{rhythm}">'
            f'{inner}<div class="pagenum">1 / 1</div></section>')


def slide(idx: int, layout: str = "S03", rhythm: str = "dense", **kw) -> dict:
    s = {"index": idx, "layout": layout, "rhythm": rhythm, "title": f"页{idx}",
         "content": kw.pop("content", {"points": [{"h": "要点", "p": "说明"}]}),
         "speaker_notes": "讲稿。"}
    if layout not in {"S01", "S02", "S07", "S08"}:
        s.update({"takeaway": "指标下降了 3 成",
                  "evidence": [{"scenario": True}, {"scenario": True}],
                  "confidence": "medium"})
    s.update(kw)
    return s


def expect(name: str, code: int, out: str, marker: str, want_fail: bool = True) -> None:
    global PASS, FAIL
    caught = code != 0 and marker in out
    ok = caught if want_fail else (code == 0)
    if ok:
        PASS += 1
        print(f"PASS  {name}")
    else:
        FAIL += 1
        print(f"FAIL  {name}: exit={code}, marker {marker!r} "
              f"{'missing' if marker not in out else 'present'}\n--- output ---\n{out}")


def deck_case(name: str, inner: str, marker: str, single: bool = True, **pagekw) -> None:
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page(inner, **pagekw)})
        args = [str(proj)] + (["--single"] if single else [])
        code, out = run("validate_deck.py", *args)
        expect(name, code, out, marker)


def deck_nofire(name: str, inner: str, marker: str, **pagekw) -> None:
    """positive control: the given marker/gate must NOT fire on legit input."""
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page(inner, **pagekw)})
        code, out = run("validate_deck.py", str(proj), "--single")
        expect(name, 1 if marker in out else 0, out, marker, want_fail=False)


def outline_case(name: str, outline: dict, marker: str, want_fail: bool = True) -> None:
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {}, outline)
        code, out = run("validate_outline.py", str(proj))
        expect(name, code, out, marker, want_fail)


def main() -> int:
    no_chrome = "--no-chrome" in sys.argv

    # ── validate_deck: every dodge the adversarial review found must now be caught
    deck_case("E1 hex color", '<div style="color:#ff0000">x</div>', "E1")
    deck_case("E1 named color", '<div style="background:crimson">x</div>', "E1")
    deck_case("E1 named color shorthand", '<div style="border:1px solid red">x</div>', "E1")
    deck_case("E1 SVG fill attr", '<svg><rect fill="#00ff00"/></svg>', "E1")
    deck_case("E1 SVG stop-color attr",
              '<svg><stop stop-color="gold"/></svg>', "E1")
    deck_case("E2 em font-size", '<div style="font-size:2.5em">x</div>', "E2")
    deck_case("E2 rem font-size", '<div style="font-size:1.2rem">x</div>', "E2")
    deck_case("E4 css gradient",
              '<div style="background:linear-gradient(#fff,#000)">x</div>', "E4")
    deck_case("E4 SVG gradient element",
              '<svg><linearGradient id="g"/></svg>', "E4")
    deck_case("E6 generator signature", '<div>Generated by deckcraft</div>', "E6")
    deck_case("E8 stray close tag", '<div>x</div></section><section class="slide">', "E8")
    deck_case("E9 display:none", '<div style="display:none">hidden</div>', "E9")
    deck_case("E9 opacity:0", '<div style="opacity:0">hidden</div>', "E9")
    # single-quote / spaced-`=` style attrs must NOT escape the style gates —
    # Chrome renders them identically, so a double-quote-only regex was a hole.
    deck_case("E1 single-quote hex", "<div style='color:#ff0000'>x</div>", "E1")
    deck_case("E2 spaced-eq font-size", '<div style = "font-size:2.5em">x</div>', "E2")
    deck_case("E9 single-quote opacity:0", "<div style='opacity:0'>hidden</div>", "E9")

    # ── false-positive guards: legit input the gates must NOT flag
    deck_nofire("E9 opacity:0.5 allowed", '<div style="opacity:0.5">dim</div>', "E9")
    # E1 must not read a color word out of a token name or an image filename
    deck_nofire("E1 var token name allowed",
                '<div style="color:var(--c-gold);background:var(--c-navy-2)">x</div>', "E1")
    deck_nofire("E1 url filename allowed",
                '<div style="background:url(images/gold-bar.png) no-repeat">x</div>', "E1")
    deck_nofire("E1 data-fill attr allowed", '<svg><rect data-fill="red" fill="currentColor"/></svg>', "E1")
    # E6 must not flag legit prose that merely contains "generated by" / "自动生成"
    deck_nofire("E6 prose generated-by allowed", '<div>营收由 Q3 生成，revenue generated by growth</div>', "E6")
    deck_nofire("E6 prose auto-gen allowed", '<div>本页讲自动生成代码的工程实践</div>', "E6")

    # E10: scenario page whose visible text lacks the 示意/虚构 label
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page("<h1>看似真实的数据页</h1>")},
                       {"title": "t", "slides": [slide(1)]})
        code, out = run("validate_deck.py", str(proj), "--single")
        expect("E10 scenario without label", code, out, "E10")
    # E10: an English deck may label in English (must NOT be forced to embed 示意)
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page("<h1>Pilot metrics</h1><small>Illustrative — not real data</small>")},
                       {"title": "t", "language": "en", "slides": [slide(1)]})
        code, out = run("validate_deck.py", str(proj), "--single")
        expect("E10 english label allowed", 1 if "E10" in out else 0, out, "E10", want_fail=False)
    # E10: English scenario deck with NO label must still fail
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page("<h1>Pilot metrics look real</h1>")},
                       {"title": "t", "language": "en", "slides": [slide(1)]})
        code, out = run("validate_deck.py", str(proj), "--single")
        expect("E10 english without label", code, out, "E10")

    # E5: page-number hole (page-1 + page-3 vs contiguous outline)
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page("<h1>甲</h1>", layout="S03"),
                            3: page("<h1>乙</h1>", layout="S04")},
                       {"title": "t", "slides": [slide(1), slide(2, layout="S04")]})
        code, out = run("validate_deck.py", str(proj))
        expect("E5 page-number hole", code, out, "E5")

    # ── validate_outline
    outline_case("O5 buzzword same-field",
                 {"title": "t", "slides": [
                     slide(1, takeaway="全面提升研发效能，打造降本增效抓手",
                           content={"points": [{"h": "01", "p": "编号在别的字段"}]})]},
                 "O5")
    outline_case("O8 title over budget",
                 {"title": "t", "slides": [
                     slide(1, title="这是一个明显超过十八个全角字符预算的超长页标题啊")]},
                 "O8")
    outline_case("O8 S03 too many points",
                 {"title": "t", "slides": [
                     slide(1, content={"points": [{"h": f"点{i}", "p": "说明"}
                                                  for i in range(6)]})]},
                 "O8")
    outline_case("O6 breathing ratio",
                 {"title": "t", "slides":
                  [slide(i, layout="S03" if i % 2 else "S04") for i in range(1, 17)]},
                 "O6")
    # positive control: a compliant 2-page outline passes
    outline_case("outline positive control",
                 {"title": "t", "slides": [
                     slide(1, layout="S01", rhythm="anchor",
                           title="短标题", content={"kicker": "K", "subtitle": "副题"}),
                     slide(2, layout="S03")]},
                 "", want_fail=False)

    # ── build_deck: dangling image reference must fail the build
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<img src="images/nope.png">')})
        code, out = run("build_deck.py", str(proj))
        expect("build missing image", code, out, "missing image")

    png_1x1 = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d49444154789c626001000000ffff03000006000557bfabd4000000"
        "0049454e44ae426082")

    # build_deck must inline ./images/ and xlink:href spellings, not just src="images/"
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page(
            '<img src="./images/a.png"><svg><image xlink:href="images/b.png"/></svg>')})
        img = proj / "images"; img.mkdir()
        (img / "a.png").write_bytes(png_1x1)
        (img / "b.png").write_bytes(png_1x1)
        code, out = run("build_deck.py", str(proj))
        deck = (proj / "deck.html").read_text(encoding="utf-8") if code == 0 else ""
        ok = code == 0 and deck.count("data:image/png;base64,") == 2 and "images/" not in deck
        expect("build inlines ./ and xlink:href", 0 if ok else 1, out if not ok else "",
               "", want_fail=False)

    # an unsupported spelling that leaves a live images/ ref must ERROR, not silently ship
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<img data-src="/abs/path/images/x.png" src="images/x.png">')})
        img = proj / "images"; img.mkdir()
        (img / "x.png").write_bytes(png_1x1)
        code, out = run("build_deck.py", str(proj))
        expect("build residual images/ ref caught", code, out, "survived inlining")

    if not no_chrome:
        # probe: an oversized image must be caught AFTER the load event (H2)
        with tempfile.TemporaryDirectory() as td:
            proj = project(td, {1: page('<img src="images/tall.png">')})
            img = proj / "images"; img.mkdir()
            # tiny valid PNG scaled up via attributes is simpler than PIL: use
            # a 1x1 PNG + explicit width/height so layout overflows physically
            png_1x1 = bytes.fromhex(
                "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
                "890000000d49444154789c626001000000ffff03000006000557bfabd4000000"
                "0049454e44ae426082")
            (img / "tall.png").write_bytes(png_1x1)
            (proj / "pages" / "page-01.html").write_text(
                page('<img src="images/tall.png" width="200" height="1600">'),
                encoding="utf-8")
            code, out = run("build_deck.py", str(proj))
            if code != 0:
                expect("probe oversized image (build)", code, out, "", want_fail=False)
            else:
                code, out = run("probe_overflow.py", str(proj))
                expect("probe oversized image", code, out, "out-of-canvas")

    if not no_chrome:
        shots_reminder_cases()

    if not no_chrome:
        editable_cases()

    print(f"\ndeckcraft-selftest: {PASS} passed · {FAIL} failed")
    return 1 if FAIL else 0


def node_available() -> bool:
    code, _ = run("find_node.py")
    return code == 0


def shots_reminder_cases() -> None:
    """--shots must hand the caller a reminder that nobody has looked at them yet.

    L4 review §四, two runs: one sent the screenshots to a sub-agent that replied
    "does not support image input" and reviewed 0 pages; the next skipped the
    review entirely. Both delivery reports listed the machine gates as all-green
    and never mentioned the visual pass. Nothing could have caught it — an
    unreviewed deck and a reviewed one are byte-identical, and the only place the
    requirement lived was SKILL.md Phase 6.

    What IS checkable is whether the reminder reaches the model's context at the
    moment it holds the screenshots. That is all this asserts. It cannot assert
    that the model then tells the user — no gate can; §二 measured the same model
    saying it on one run and not on the next.
    """
    global PASS, FAIL
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<h1>x</h1>')})
        code, out = run("build_deck.py", str(proj))
        if code != 0:
            expect("shots reminder (build)", code, out, "", want_fail=False)
            return
        code, out = run("export_deck.py", str(proj), "--shots")
        if code != 0:
            expect("shots reminder (export)", code, out, "", want_fail=False)
            return
        for needle in ("NEXT:", "视觉审查", "pages_reviewed"):
            expect(f"--shots reminder names {needle!r}",
                   0 if needle in out else 1, out, needle, want_fail=False)

        # Negative control: the same export with the reminder removed — i.e. the
        # implementation that shipped until 2026-08-16. Without this, the three
        # assertions above are also satisfied by a run that prints the reminder
        # for some unrelated reason.
    with tempfile.TemporaryDirectory() as td2:
        stripped = Path(td2) / "scripts-silent"
        shutil.copytree(SCRIPTS, stripped)
        target = stripped / "export_deck.py"
        text = target.read_text(encoding="utf-8")
        anchor = '        print(f"NEXT: 这 {n_pages} 张截图还没有任何人看过。'
        if text.count(anchor) != 1:
            FAIL += 1
            print(f"FAIL  --shots reminder control: anchor matched "
                  f"{text.count(anchor)} times, expected 1 — the control did not "
                  f"replicate the pre-2026-08-16 implementation")
            return
        head, _, tail = text.partition(anchor)
        # drop the reminder block: from the anchor to the next top-level `if`
        rest = tail.split("\n    if a.pptx:", 1)
        if len(rest) != 2:
            FAIL += 1
            print("FAIL  --shots reminder control: could not find the end of the "
                  "reminder block — the control did not apply")
            return
        target.write_text(head + "\n    if a.pptx:" + rest[1], encoding="utf-8")
        proj2 = project(td2, {1: page('<h1>x</h1>')})
        subprocess.run([PY, str(SCRIPTS / "build_deck.py"), str(proj2)],
                       capture_output=True, text=True, timeout=180)
        r = subprocess.run([PY, str(target), str(proj2), "--shots"],
                           capture_output=True, text=True, encoding="utf-8",
                           errors="replace", timeout=180)
        silent = (r.stdout or "") + (r.stderr or "")
        if r.returncode != 0:
            FAIL += 1
            print(f"FAIL  --shots reminder control: the stripped export itself "
                  f"failed (exit {r.returncode}) — a control that cannot run is "
                  f"not a control\n--- output ---\n{silent}")
            return
        fired = [n for n in ("NEXT:", "视觉审查", "pages_reviewed") if n in silent]
        expect("--shots reminder control: silent build says none of it",
               1 if fired else 0, f"still present: {fired}", "", want_fail=False)


def editable_cases() -> None:
    """P2b --pptx-editable: text/shape translation, raster degrade, no leak.
    Node-gated (skipped with a note when no runtime is present — the export path
    fails fast with guidance in that case, which we assert separately)."""
    global PASS, FAIL
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if not node_available():
        # The export path fails fast with install guidance when Node is absent
        # (asserted here); the translation cases below need a real runtime.
        with tempfile.TemporaryDirectory() as td:
            proj = project(td, {1: page('<h1>x</h1>')})
            run("build_deck.py", str(proj))
            code, out = run("export_deck.py", str(proj), "--pptx-editable")
            expect("editable no-node guided fail", code, out, "Node.js runtime")
        print("SKIP  editable-pptx translation cases (no Node runtime)")
        return

    def build_editable(td: str, inner: str) -> "Presentation | None":
        proj = project(td, {1: page(inner)})
        code, out = run("build_deck.py", str(proj))
        if code != 0:
            expect("editable build_deck", code, out, "", want_fail=False)
            return None
        code, out = run("export_deck.py", str(proj), "--pptx-editable")
        pptx_f = proj / "export" / "deck-editable.pptx"
        if code != 0 or not pptx_f.is_file():
            expect("editable export", 1, out, "", want_fail=False)
            return None
        # nothing may leak outside export/ — intermediates stay in the dot-safe export dir
        strays = [p.name for p in proj.iterdir()
                  if p.name not in {"pages", "tokens.css", "export", "deck.html", "outline.json"}]
        expect("editable no stray files", 1 if strays else 0,
               f"strays: {strays}", "", want_fail=False)
        return Presentation(str(pptx_f))

    # 1) text + shape → editable text box (exact text) + a shape, positioned right
    with tempfile.TemporaryDirectory() as td:
        prs = build_editable(td,
            '<div style="position:absolute;left:120px;top:200px;width:400px">'
            '可编辑正文 Hello</div>'
            '<div style="position:absolute;left:640px;top:120px;width:200px;'
            'height:150px;background:var(--c-primary)"></div>')
        if prs:
            s = prs.slides[0]
            from pptx.util import Emu
            txts = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
            hit = next((sh for sh in txts if "Hello" in sh.text_frame.text), None)
            expect("editable text box present", 0 if hit else 1, "", "", want_fail=False)
            if hit:
                x_in = Emu(hit.left).inches
                # left:120px → 120/96 = 1.25in, tolerate metric/padding drift
                expect("editable text positioned", 0 if abs(x_in - 1.25) < 0.4 else 1,
                       f"x={x_in:.2f}in (want ~1.25)", "", want_fail=False)
            shapes = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                      and not (sh.has_text_frame and sh.text_frame.text.strip())]
            expect("editable shape present", 0 if shapes else 1, "", "", want_fail=False)

    # 2) inline SVG → honest raster: a picture + the delivery NOTE, text still editable
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page(
            '<p>可编辑文字</p>'
            '<svg viewBox="0 0 24 24" width="80" height="80" stroke="currentColor" '
            'fill="none"><circle cx="12" cy="12" r="9"/></svg>')})
        code, _ = run("build_deck.py", str(proj))
        code, out = run("export_deck.py", str(proj), "--pptx-editable")
        expect("raster degrade reported", 0 if ("栅格化" in out and code == 0) else 1,
               out, "", want_fail=False)
        pptx_f = proj / "export" / "deck-editable.pptx"
        if pptx_f.is_file():
            s = Presentation(str(pptx_f)).slides[0]
            pics = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            txts = [sh for sh in s.shapes if sh.has_text_frame and sh.text_frame.text.strip()]
            expect("raster produced a picture", 0 if pics else 1, "", "", want_fail=False)
            expect("raster kept text editable", 0 if txts else 1, "", "", want_fail=False)

    # 3) layout.json schema shape
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<h1>甲</h1><p>乙</p>')})
        run("build_deck.py", str(proj))
        code, out = run("extract_layout.py", str(proj))
        ok = False
        if code == 0:
            try:
                d = json.loads(out)
                pg = d["pages"][0]
                ok = (d["canvas"]["w"] == 1280 and pg["index"] == 1
                      and all("type" in e and "box" in e for e in pg["elements"])
                      and any(e["type"] == "text" for e in pg["elements"]))
            except (json.JSONDecodeError, KeyError, IndexError):
                ok = False
        expect("layout.json schema", 0 if ok else 1, out, "", want_fail=False)

    # 3a) a literal </script> in visible text must NOT break the <script>JSON</script>
    #     extraction — guarded by the injection-side '<' → < escape (extract + probe).
    #     Without it, --dump-dom would close the JSON block early → json.loads crash.
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<h1>Web</h1><p>code sample: &lt;/script&gt; here</p>')})
        run("build_deck.py", str(proj))
        code, out = run("extract_layout.py", str(proj))
        ok = False
        if code == 0:
            try:
                ok = "</script>" in json.dumps(json.loads(out))
            except json.JSONDecodeError:
                ok = False
        expect("extract survives </script> in visible text", 0 if ok else 1, out, "", want_fail=False)

    # 3b) presenter-view parity: the editable pptx must carry outline speaker_notes
    #     (the image-type pptx does; choosing editable must not silently drop them)
    with tempfile.TemporaryDirectory() as td:
        proj = project(td, {1: page('<h1>标题</h1><p>正文</p>')},
                       {"title": "t", "slides": [slide(1)]})
        run("build_deck.py", str(proj))
        run("export_deck.py", str(proj), "--pptx-editable")
        pptx_f = proj / "export" / "deck-editable.pptx"
        has = False
        if pptx_f.is_file():
            sl = Presentation(str(pptx_f)).slides[0]
            has = sl.has_notes_slide and "讲稿" in sl.notes_slide.notes_text_frame.text
        expect("editable carries speaker notes", 0 if has else 1, "", "", want_fail=False)

    # 4) classification at the extract seam (adversarial-review HIGH-2 + MED-3):
    #    a non-data <img> must become a RASTER (never a path the assembler can't
    #    resolve → silent drop); a data: <img> stays an embeddable image; a
    #    background-image container reports the editable text it swallows (textLost).
    px = ("data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwC"
          "AAAAC0lEQVR42mNk+M8AAAMBAQDJ/pLvAAAAAElFTkSuQmCC")
    with tempfile.TemporaryDirectory() as td:
        proj = Path(td) / "proj"
        proj.mkdir()
        (proj / "deck.html").write_text(
            '<html><head><style>.slide{width:1280px;height:720px;position:relative}'
            '</style></head><body><div class="stage">'
            '<section class="slide" data-layout="S03" data-rhythm="dense">'
            f'<div style="background-image:url({px});width:300px;height:200px">'
            '<h3>卡标题</h3><span>卡文字</span></div>'
            '<img src="pics/logo.png" width="120" height="60">'
            f'<img src="{px}" width="40" height="40">'
            '</section></div></body></html>', encoding="utf-8")
        code, out = run("extract_layout.py", str(proj))
        ok = False
        if code == 0:
            try:
                els = json.loads(out)["pages"][0]["elements"]
                rasters = [e for e in els if e["type"] == "raster"]
                images = [e for e in els if e["type"] == "image"]
                bg = next((e for e in rasters if e.get("tag") == "bg-image"), None)
                ok = (any(e.get("tag") == "img" for e in rasters)  # relative <img> → raster
                      and len(images) == 1                          # data: <img> → image
                      and bg is not None and bg.get("textLost") == 2)  # swallowed text counted
            except (json.JSONDecodeError, KeyError, IndexError):
                ok = False
        expect("img/bg-image raster classification + textLost", 0 if ok else 1, out, "", want_fail=False)

    # 5) fidelity at the extract seam (adversarial-review H1/M1/M5): inline caption
    #    span keeps its own smaller size + opacity; a single-side border becomes a
    #    thin rect (not dropped); a padded text leaf is inset to its content box.
    with tempfile.TemporaryDirectory() as td:
        proj = Path(td) / "proj"
        proj.mkdir()
        (proj / "deck.html").write_text(
            '<html><head><style>.slide{width:1280px;height:720px;position:relative}'
            '</style></head><body><div class="stage">'
            '<section class="slide" data-layout="S03" data-rhythm="dense">'
            '<p style="position:absolute;left:96px;top:96px;font-size:24px">大<span '
            'style="font-size:12px;opacity:.5">小</span></p>'
            '<div style="position:absolute;left:96px;top:240px;width:400px;height:100px;'
            'border-top:8px solid #C75B12"></div>'
            '<div style="position:absolute;left:600px;top:240px;width:300px;font-size:18px;'
            'padding-left:32px;color:#111111">缩进</div>'
            '</section></div></body></html>', encoding="utf-8")
        code, out = run("extract_layout.py", str(proj))
        ok = False
        if code == 0:
            try:
                els = json.loads(out)["pages"][0]["elements"]
                big = next(e for e in els if e["type"] == "text"
                           and any("大" in r.get("text", "") for r in e["runs"]))
                cap = next(r for r in big["runs"] if "小" in r.get("text", ""))
                stripe = [e for e in els if e["type"] == "rect"
                          and e.get("fill") == "C75B12" and e["box"]["h"] == 8]
                pad = next(e for e in els if e["type"] == "text"
                           and any("缩进" in r.get("text", "") for r in e["runs"]))
                ok = (cap["fontPx"] == 12 and abs(cap["opacity"] - 0.5) < 0.01      # H1
                      and len(stripe) == 1                                          # M1
                      and pad["box"]["x"] == 632)                                   # M5 (600+32)
            except (json.JSONDecodeError, KeyError, IndexError, StopIteration):
                ok = False
        expect("per-run font/opacity + single-side border + padding inset",
               0 if ok else 1, out, "", want_fail=False)

    # 6) no spurious re-wrap (real-machine finding): a browser-single-line number/
    #    label must be marked wrap:false (PowerPoint's wider font would otherwise
    #    stack "01"→"0"/"1"); a genuinely multi-line paragraph keeps wrap:true.
    with tempfile.TemporaryDirectory() as td:
        proj = Path(td) / "proj"
        proj.mkdir()
        (proj / "deck.html").write_text(
            '<html><head><style>.slide{width:1280px;height:720px;position:relative}'
            '.para{width:360px;font-size:20px;line-height:1.6}</style></head>'
            '<body><div class="stage">'
            '<section class="slide" data-layout="S02" data-rhythm="dense">'
            '<div style="position:absolute;left:760px;top:230px;font-size:40px">01</div>'
            '<div class="para" style="position:absolute;left:96px;top:400px">这是一段明显'
            '需要换行的正文，它在浏览器里就跨了多行，pptx 也应保持换行不能挤成一行。</div>'
            '</section></div></body></html>', encoding="utf-8")
        code, out = run("extract_layout.py", str(proj))
        ok = False
        if code == 0:
            try:
                els = json.loads(out)["pages"][0]["elements"]
                num = next(e for e in els if e["type"] == "text"
                           and any(r.get("text") == "01" for r in e["runs"]))
                para = next(e for e in els if e["type"] == "text"
                            and any("这是一段" in r.get("text", "") for r in e["runs"]))
                ok = num.get("wrap") is False and para.get("wrap") is True
            except (json.JSONDecodeError, KeyError, IndexError, StopIteration):
                ok = False
        expect("single-line number wrap:false, multi-line para wrap:true",
               0 if ok else 1, out, "", want_fail=False)


if __name__ == "__main__":
    sys.exit(main())
