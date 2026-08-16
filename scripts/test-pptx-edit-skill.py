#!/usr/bin/env python3
"""Behaviour tests for the `pptx-edit` skill (discussions/059 §六·补八).

    python3 scripts/test-pptx-edit-skill.py
    python3 scripts/test-pptx-edit-skill.py --json

WHY THIS EXISTS
---------------
Two of these scripts' properties are LIMITS, not features — table/group text is
invisible to both of them, and a phrase PowerPoint split across runs will not
match. S6 wrote those limits into SKILL.md as measured facts. A claim in a
SKILL.md that nothing checks is a claim that rots: the day someone teaches
pptx_read.py to walk tables, SKILL.md starts lying in the other direction.
So the limits are asserted here, in both directions.

The other half is a regression guard. The S6 wrap-up review fed the scripts a
.pptx built by a minimal OOXML generator rather than by PowerPoint — no
slideLayout relationship, no slideMaster — and BOTH scripts died with a bare
traceback (KeyError in pptx_read, IndexError in pptx_edit's own bounds check).
That is the failure mode the whole family is held to: adversarial input must
exit non-zero with one sentence, never a traceback.

A third defect only exists on Windows: neither script reconfigured stdout to
UTF-8, so a CAPTURED stdout — which is how an agent always calls them — died with
UnicodeEncodeError on the first Chinese character. It shipped that way from the
doc-edit days and was found by this gate's FIRST CI run. C1 reproduces the
condition portably by forcing an ANSI code page on the child, so it is now
catchable on macOS and Linux too.

Every assertion runs twice: once against the real scripts (must stay silent) and
once against a state carrying exactly the defect it hunts (must fire). The five
LIVE controls do not fabricate output — they re-run the REAL scripts with the
real fix reverted, i.e. they replicate the implementation that actually shipped
until 2026-08-04. A control that cannot tell the two implementations apart is
not a control — and the UTF-8 one is the sharpest example: on a UTF-8 machine the
guarded and unguarded scripts are byte-for-byte identical in behaviour.

Lives outside skills/builtin/ so it is not packed into skills-builtin.zip.
Exit 0 = every assertion behaved, 1 = something did not.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

# stdout must be UTF-8 on every platform, and on Windows a captured stdout is not
# (same defect the other gates in this directory were fixed for).
for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        try:
            _stream.reconfigure(encoding="utf-8")
        except (ValueError, OSError):
            pass

REPO = Path(__file__).resolve().parent.parent
SKILL = REPO / "skills" / "builtin" / "pptx-edit"
SCRIPTS = SKILL / "scripts"
PY = sys.executable

# The probe appears TWICE in the fixture: once in a plain textbox and once in a
# table cell. That is what makes the silent-miss assertion possible — the tool
# replaces one, leaves the other, and reports a count that mentions neither fact.
PROBE = "营业收入同比增长"
TABLE_ONLY = "表内独有文字"
GROUP_ONLY = "组合内独有文字"
SPLIT_HEAD, SPLIT_TAIL = "毛利", "率保持稳定"   # one phrase, two runs
SINGLE_RUN = "费用率同比下降"
ABSENT = "这份文件里根本没有的词"               # matches nothing, anywhere


# ── fixtures ──────────────────────────────────────────────────────────────────
def build_deck(path: Path, *, with_table=True, split_phrase=True, with_group=True) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "季度经营分析"

    tb = s.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(4), Inches(0.8))
    tb.text_frame.text = PROBE

    tb2 = s.shapes.add_textbox(Inches(0.4), Inches(2.6), Inches(4), Inches(0.8))
    p = tb2.text_frame.paragraphs[0]
    if split_phrase:
        for chunk in (SPLIT_HEAD, SPLIT_TAIL):
            p.add_run().text = chunk
    else:
        p.add_run().text = SPLIT_HEAD + SPLIT_TAIL        # one run: control arm

    tb3 = s.shapes.add_textbox(Inches(0.4), Inches(3.6), Inches(4), Inches(0.8))
    tb3.text_frame.text = SINGLE_RUN + " 1.2 个百分点"

    if with_table:
        t = s.shapes.add_table(2, 2, Inches(5), Inches(1.6), Inches(4), Inches(1.2)).table
        t.cell(0, 0).text = TABLE_ONLY
        t.cell(1, 1).text = PROBE                          # the second copy

    if with_group:
        g = s.shapes.add_group_shape()
        gb = g.shapes.add_textbox(Inches(5), Inches(3.6), Inches(3), Inches(0.8))
        gb.text_frame.text = GROUP_ONLY

    prs.save(str(path))


def build_layoutless(src: Path, dst: Path) -> None:
    """A .pptx that opens but carries no slideLayout/slideMaster.

    Minimal OOXML generators produce these; PowerPoint does not. Both crashes the
    S6 review found were on a file shaped exactly like this.
    """
    drop = re.compile(r"^ppt/(slideLayouts|slideMasters|theme)/")
    strip_rel = re.compile(r"<Relationship[^>]*(?:slideLayout|slideMaster|theme)[^>]*/>")
    strip_ovr = re.compile(r"<Override[^>]*(?:slideLayout|slideMaster|theme)[^>]*/>")
    strip_lst = re.compile(r"<p:sldMasterIdLst>.*?</p:sldMasterIdLst>", re.S)
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zo:
        for item in zin.infolist():
            if drop.match(item.filename):
                continue
            data = zin.read(item.filename)
            if item.filename.endswith(".rels") or item.filename == "[Content_Types].xml":
                text = strip_ovr.sub("", strip_rel.sub("", data.decode("utf-8")))
                data = text.encode("utf-8")
            if item.filename == "ppt/presentation.xml":
                data = strip_lst.sub("", data.decode("utf-8")).encode("utf-8")
            zo.writestr(item, data)


def build_adversarial(work: Path) -> dict[str, Path]:
    out = {}
    (work / "empty.pptx").write_bytes(b"")
    out["0 字节"] = work / "empty.pptx"
    (work / "notzip.pptx").write_bytes(b"this is plainly not a zip")
    out["非 zip"] = work / "notzip.pptx"
    (work / "adir.pptx").mkdir()
    out["目录"] = work / "adir.pptx"
    with zipfile.ZipFile(work / "emptyzip.pptx", "w") as z:
        z.writestr("hello.txt", "hi")
    out["zip 但不是 pptx"] = work / "emptyzip.pptx"
    out["不存在"] = work / "missing.pptx"
    return out


# ── running the scripts ───────────────────────────────────────────────────────
def run(script_dir: Path, name: str, *args: str, env: dict | None = None) -> dict:
    proc = subprocess.run([PY, str(script_dir / name), *map(str, args)],
                          capture_output=True, text=True, encoding="utf-8",
                          env={**os.environ, **env} if env else None)
    return {"exit": proc.returncode, "out": proc.stdout or "", "err": proc.stderr or ""}


# An ANSI code page with no room for CJK, forced on the child. This is the same
# condition a Windows agent creates simply by CAPTURING the script's stdout — the
# runner's code page, not the terminal's, is what Python encodes to. Reproducing it
# through PYTHONIOENCODING is what makes the assertion runnable on macOS and Linux
# too; without that, the defect is invisible everywhere except a Windows CI run
# (which is exactly how it survived from the doc-edit days until 2026-08-04).
ANSI_ENV = {"PYTHONIOENCODING": "cp1252"}


def default_captured_encoding() -> str:
    """What a child process ACTUALLY gets for stdout when its output is captured.

    Measured, not assumed: on macOS/Linux this is utf-8, on Windows it is the
    machine's ANSI code page. It decides how much of the skill breaks when the
    UTF-8 reconfigure is removed — on a UTF-8 host only the forced-ANSI assertion
    (C1) can see the difference, while on Windows EVERY run of the script breaks,
    because there the default already is the hostile code page. Declaring one
    fixed cascade for that control would be wrong on one of the two platforms;
    the first CI run of this gate proved it by going red on Windows alone.
    """
    proc = subprocess.run([PY, "-c", "import sys; print(sys.stdout.encoding)"],
                          capture_output=True, text=True)
    return (proc.stdout or "").strip().lower()


CAPTURED_ENC = default_captured_encoding()
HOST_CAPTURES_UTF8 = CAPTURED_ENC.replace("-", "") == "utf8"


def sha(p: Path) -> str:
    return hashlib.sha256(p.read_bytes()).hexdigest()


def collect(work: Path, script_dir: Path, *, with_table=True, split_phrase=True,
            with_group=True, layoutless_has_layouts=False) -> dict:
    work.mkdir(parents=True, exist_ok=True)
    deck = work / "deck.pptx"
    build_deck(deck, with_table=with_table, split_phrase=split_phrase, with_group=with_group)

    nolayout = work / "nolayout.pptx"
    if layoutless_has_layouts:
        shutil.copyfile(deck, nolayout)      # control arm: it is NOT layout-less
    else:
        build_layoutless(deck, nolayout)

    ctx: dict = {"work": work, "script_dir": script_dir}
    ctx["read"] = run(script_dir, "pptx_read.py", deck)
    # the same read, but with a CJK-hostile code page forced on the child
    ctx["read_ansi"] = run(script_dir, "pptx_read.py", deck, env=ANSI_ENV)
    ctx["read_json"] = run(script_dir, "pptx_read.py", deck, "--json")
    ctx["read_nolayout"] = run(script_dir, "pptx_read.py", nolayout)

    # Every edit gets its OWN copy of the deck. Sharing one input makes the
    # assertions order-dependent: an implementation that ignores --out would edit
    # the shared file in place and every later measurement would be reading a
    # document the previous case already changed.
    def fresh(tag: str) -> Path:
        cp = work / f"in-{tag}.pptx"
        shutil.copyfile(deck, cp)
        return cp

    # replace the probe, which lives in a textbox AND in a table cell
    src1, out1 = fresh("replace"), work / "replaced.pptx"
    before = sha(src1)
    ctx["replace"] = run(script_dir, "pptx_edit.py", src1, "--replace", PROBE, "XX",
                         "--out", out1)
    ctx["replace_ansi"] = run(script_dir, "pptx_edit.py", fresh("ansi"), "--replace",
                              PROBE, "XX", "--out", work / "ansi.pptx", env=ANSI_ENV)
    ctx["input_untouched"] = sha(src1) == before
    ctx["surviving_probe"] = count_everywhere(out1, PROBE) if out1.exists() else None
    ctx["replaced_textbox"] = count_textframes(out1, "XX") if out1.exists() else None

    out2 = work / "split.pptx"
    ctx["replace_split"] = run(script_dir, "pptx_edit.py", fresh("split"), "--replace",
                               SPLIT_HEAD + SPLIT_TAIL, "YY", "--out", out2)
    out3 = work / "single.pptx"
    ctx["replace_single"] = run(script_dir, "pptx_edit.py", fresh("single"), "--replace",
                                SINGLE_RUN, "ZZ", "--out", out3)
    ctx["single_applied"] = count_textframes(out3, "ZZ") if out3.exists() else None

    # A string that exists ONLY inside the group. The table case alone cannot tell
    # a census that walks tables from one that walks both.
    ctx["replace_group"] = run(script_dir, "pptx_edit.py", fresh("group"), "--replace",
                               GROUP_ONLY, "GG", "--out", work / "group.pptx")

    # ── in-place writes: does a run that changed nothing still rewrite the file? ──
    # No --out on purpose. Every case above passes --out, so the in-place path — the
    # DEFAULT one, and the one L4 §四 caught the model taking on the user's own deck
    # — had no coverage at all until now.
    noop = fresh("noop")
    noop_before = sha(noop)
    ctx["replace_noop"] = run(script_dir, "pptx_edit.py", noop, "--replace", ABSENT, "QQ")
    ctx["noop_untouched"] = sha(noop) == noop_before

    # The same shape, but the needle exists ONLY in a table: 0 replacements AND
    # something out of reach. This is 🚧C5 verbatim — E in the L4 §四 rerun ran
    # exactly this against the real deck.
    unreach = fresh("unreach")
    unreach_before = sha(unreach)
    ctx["replace_unreachable"] = run(script_dir, "pptx_edit.py", unreach,
                                     "--replace", TABLE_ONLY, "QQ")
    ctx["unreachable_untouched"] = sha(unreach) == unreach_before

    # The other direction. Without this, an implementation that never writes in
    # place at all satisfies both cases above — and "it silently did nothing" is
    # a worse defect than the one being fixed.
    inplace = fresh("inplace")
    inplace_before = sha(inplace)
    ctx["replace_inplace"] = run(script_dir, "pptx_edit.py", inplace,
                                 "--replace", SINGLE_RUN, "WW")
    ctx["inplace_written"] = sha(inplace) != inplace_before
    ctx["inplace_applied"] = count_textframes(inplace, "WW")

    ctx["add_nolayout"] = run(script_dir, "pptx_edit.py", nolayout, "--add-slide",
                              "--layout", "0", "--out", work / "added.pptx")
    ctx["layout_high"] = run(script_dir, "pptx_edit.py", fresh("hi"), "--add-slide",
                             "--layout", "99", "--out", work / "hi.pptx")
    ctx["layout_neg"] = run(script_dir, "pptx_edit.py", fresh("neg"), "--add-slide",
                            "--layout", "-1", "--out", work / "neg.pptx")

    adv = build_adversarial(work)
    ctx["adversarial"] = {}
    for label, path in adv.items():
        ctx["adversarial"][label] = {
            "read": run(script_dir, "pptx_read.py", path),
            "edit": run(script_dir, "pptx_edit.py", path, "--replace", "a", "b",
                        "--out", work / "adv-out.pptx"),
        }

    # what the fixture actually contains, read back with the library directly —
    # never trust the tool under test to describe its own input
    ctx["fixture"] = describe(deck)
    ctx["nolayout_really_raises"] = layoutless_raises(nolayout)
    return ctx


def describe(path: Path) -> dict:
    from pptx import Presentation
    prs = Presentation(str(path))
    tf_texts, cell_texts, runs, group_texts = [], [], [], []
    n_tables = n_groups = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if getattr(sh, "has_table", False) and sh.has_table:
                n_tables += 1
            elif getattr(sh, "shapes", None) is not None and not sh.has_text_frame:
                n_groups += 1
            if sh.has_text_frame:
                tf_texts.append(sh.text_frame.text)
                for para in sh.text_frame.paragraphs:
                    if len(para.runs) > 1:
                        runs.append([r.text for r in para.runs])
            if getattr(sh, "has_table", False) and sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        cell_texts.append(cell.text)
            # A group is neither a text frame nor a table; reaching inside it is
            # the only way to prove the fixture really carries group text.
            if getattr(sh, "shapes", None) is not None and not sh.has_text_frame:
                for inner in sh.shapes:
                    if inner.has_text_frame and inner.text_frame.text.strip():
                        group_texts.append(inner.text_frame.text)
    return {"textframes": tf_texts, "cells": cell_texts, "multirun": runs,
            "group_texts": group_texts,
            "n_tables": n_tables, "n_groups": n_groups}


def unreachable_truth(path: Path, needle: str) -> int:
    """How many RUNS carrying `needle` sit where --replace cannot reach.

    Counted here with the library directly, never with the script under test —
    otherwise the assertion would be comparing the tool's report against the tool.
    Same unit the script prints: one per matching run.
    """
    from pptx import Presentation

    def tf_hits(tf) -> int:
        return sum(1 for para in tf.paragraphs for run in para.runs if needle in run.text)

    def walk(shape) -> int:
        n = 0
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    n += tf_hits(cell.text_frame)
            return n
        subs = getattr(shape, "shapes", None)
        if subs is not None and not shape.has_text_frame:
            for inner in subs:
                if inner.has_text_frame:
                    n += tf_hits(inner.text_frame)
                n += walk(inner)
        return n

    prs = Presentation(str(path))
    return sum(walk(sh) for slide in prs.slides for sh in slide.shapes)


def reported_unread(out: str) -> tuple[int, int]:
    """(tables, groups) as the READ script claims them. Absent line == (0, 0),
    which is precisely what the silent implementation reports."""
    t = re.search(r"表格 ×(\d+)", out)
    g = re.search(r"组合 ×(\d+)", out)
    return (int(t.group(1)) if t else 0, int(g.group(1)) if g else 0)


def reported_missed(out: str, needle: str) -> int:
    """What the EDIT script claims it could see but not reach. Absent == 0."""
    m = re.search(rf"另有 (\d+) 处「{re.escape(needle)}」", out)
    return int(m.group(1)) if m else 0


def count_everywhere(path: Path, needle: str) -> int:
    """Occurrences across text frames AND table cells — the full truth."""
    d = describe(path)
    return sum(t.count(needle) for t in d["textframes"] + d["cells"])


def count_textframes(path: Path, needle: str) -> int:
    return sum(t.count(needle) for t in describe(path)["textframes"])


def layoutless_raises(path: Path) -> bool:
    from pptx import Presentation
    prs = Presentation(str(path))
    try:
        _ = prs.slides[0].slide_layout.name
        return False
    except Exception:  # noqa: BLE001
        return True


def reported_count(res: dict) -> int | None:
    m = re.search(r"replacements:\s*(\d+)", res["out"])
    return int(m.group(1)) if m else None


# ── the assertions ────────────────────────────────────────────────────────────
CHECKS: dict[str, dict] = {}


def check(cid: str, title: str):
    def deco(fn):
        # A duplicate id silently REPLACES the earlier assertion while the printed
        # total keeps climbing — the xlsx gate lost N9 that way on 2026-08-16 and
        # the count never noticed. Refuse loudly instead.
        if cid in CHECKS:
            raise SystemExit(f"duplicate check id {cid!r}: the earlier assertion "
                             f"({CHECKS[cid]['title']!r}) would be silently replaced")
        CHECKS[cid] = {"id": cid, "title": title, "fn": fn}
        return fn
    return deco


@check("V0", "the fixture actually gives every assertion something to look at")
def v0(ctx: dict) -> list[str]:
    """A silent check and a check with no subject look identical from outside."""
    out = []
    f = ctx["fixture"]
    if PROBE not in f["cells"]:
        out.append("no table cell carries the probe — L1/L4 have no subject")
    if TABLE_ONLY not in f["cells"]:
        out.append("the table-only string is missing — L1 has no subject")
    if f["group_texts"] != [GROUP_ONLY]:
        out.append(f"the group shape does not carry its probe ({f['group_texts']}) "
                   f"— the group half of L1 has no subject")
    if not any(r == [SPLIT_HEAD, SPLIT_TAIL] for r in f["multirun"]):
        out.append("the probe phrase is not split across two runs — L3 has no subject")
    if sum(t.count(PROBE) for t in f["textframes"]) != 1:
        out.append("the probe does not appear exactly once in a text frame — L4 is ambiguous")
    if not ctx["nolayout_really_raises"]:
        out.append("the layout-less fixture still resolves a layout — X1/X2 have no subject")
    if ctx["read"]["exit"] != 0 or not ctx["read"]["out"].strip():
        out.append("reading the normal deck produced nothing")
    return out


@check("L1", "table and group text are invisible to pptx_read (a documented LIMIT)")
def l1(ctx: dict) -> list[str]:
    """SKILL.md 「限制」 says so in words; this is the same claim as an assertion.
    It fires in BOTH directions — if someone teaches the script to walk tables,
    this goes red and SKILL.md has to be corrected rather than silently rot."""
    out = []
    if TABLE_ONLY in ctx["read"]["out"]:
        out.append("table text now IS reported — SKILL.md 「限制」 is out of date")
    if GROUP_ONLY in ctx["read"]["out"]:
        out.append("group text now IS reported — SKILL.md 「限制」 is out of date")
    return out


@check("L3", "a phrase split across runs does not match (a documented LIMIT)")
def l3(ctx: dict) -> list[str]:
    n = reported_count(ctx["replace_split"])
    if n != 0:
        return [f"cross-run replacement reported {n}, expected 0 "
                f"(SKILL.md says it cannot match)"]
    return []


@check("L4", "a missed replacement stays missed — the reach did not change")
def l4(ctx: dict) -> list[str]:
    """The probe is in a textbox AND in a table cell. The tool changes one and
    leaves the other; `replacements:` counts the one it changed.

    This pins the REACH, which is still deliberately narrow — the skill is thin on
    purpose and the table copy must survive. What changed on 2026-08-16 is only
    whether the tool ADMITS the miss; N2 asserts that half. Keeping them apart
    matters: a future edit that quietly teaches --replace to walk tables would make
    N2 pass for the wrong reason, and only this check would catch it."""
    out = []
    if ctx["replaced_textbox"] != 1:
        out.append(f"the text-frame copy was not replaced ({ctx['replaced_textbox']})")
    if ctx["surviving_probe"] != 1:
        out.append(f"expected exactly 1 surviving copy in the table, got "
                   f"{ctx['surviving_probe']}")
    if reported_count(ctx["replace"]) != 1:
        out.append(f"reported count was {reported_count(ctx['replace'])}, expected 1")
    return out


@check("N1", "pptx_read reports HOW MANY containers it could not read, truthfully")
def n1(ctx: dict) -> list[str]:
    """L4 review §四: a model recovered the table on a real deck only because that
    slide's title happened to say 「（表格）」. Strip that hint and the page reads
    exactly like a title-only page — an empty walk looked identical to a pass.

    Asserted against the fixture's real shape count (measured with python-pptx
    here, not with the script under test), so it fires when the script goes silent
    AND when it over-reports. It deliberately does NOT check that any table TEXT
    appears — L1 asserts the opposite of that, and both must hold at once:
    the existence is reported, the content still is not.
    """
    out = []
    f = ctx["fixture"]
    want = (f["n_tables"], f["n_groups"])
    got = reported_unread(ctx["read"]["out"])
    if got != want:
        out.append(f"plain output claims {got} unread (table, group), fixture has {want}")
    try:
        payload = json.loads(ctx["read_json"]["out"])
        jt = payload.get("unread_total", {})
        gotj = (jt.get("table", 0), jt.get("group", 0))
        if gotj != want:
            out.append(f"--json unread_total {gotj}, fixture has {want}")
        if any("unread" not in s for s in payload.get("slides", [])):
            out.append("--json: a slide carries no `unread` key")
    except (ValueError, TypeError) as exc:
        out.append(f"--json output did not parse ({exc})")
    # Line-anchored on purpose. A substring test passes on the closing summary,
    # whose own wording contains the token 「上面已按页用 [unread] 标出」 — so
    # `"[unread]" in out` stayed true even with every per-slide marker suppressed.
    # Caught by the silent-read control, which was supposed to break this and did
    # not: the check was measuring the prose, not the markers.
    if want != (0, 0) and not re.search(r"(?m)^\[unread\] ", ctx["read"]["out"]):
        out.append("no per-slide [unread] marker — the caller cannot tell WHICH page")
    return out


@check("N2", "pptx_edit reports how many matches it could see but not reach")
def n2(ctx: dict) -> list[str]:
    """The other half of L4. `replacements: N` can only go up, so on its own it
    reads the same whether the run missed nothing or missed nine — 🚧C5 in the L4
    review was built to catch exactly that and could not, because the model walked
    around the script entirely.

    Truth is counted with python-pptx directly. Both the table case and the
    group case are checked: a census that walks tables but not groups passes the
    first and is still half blind.
    """
    out = []
    src = ctx["work"] / "in-replace.pptx"
    if src.exists():
        want = unreachable_truth(src, PROBE)
        got = reported_missed(ctx["replace"]["out"], PROBE)
        if got != want:
            out.append(f"table case: reported {got} unreachable, truth is {want}")
    grp = ctx["work"] / "in-group.pptx"
    if grp.exists():
        wantg = unreachable_truth(grp, GROUP_ONLY)
        gotg = reported_missed(ctx["replace_group"]["out"], GROUP_ONLY)
        if gotg != wantg:
            out.append(f"group case: reported {gotg} unreachable, truth is {wantg}")
    return out


@check("N3", "it stays quiet when there is genuinely nothing out of reach")
def n3(ctx: dict) -> list[str]:
    """A warning that fires on every run is furniture, and gets read as furniture.
    SINGLE_RUN lives in a plain textbox and nowhere else, so this replacement has
    nothing unreachable and must print no `[!]` line at all."""
    if reported_missed(ctx["replace_single"]["out"], SINGLE_RUN) != 0:
        return ["warned about unreachable matches when there were none"]
    if "[!]" in ctx["replace_single"]["out"]:
        return ["printed a [!] line on a replacement with nothing out of reach"]
    return []


@check("N4", "a run that changed nothing leaves the input file alone")
def n4(ctx: dict) -> list[str]:
    """L4 §四, measured on the user's real deck: `--replace` on a word that lives
    only in a table reported 0 replacements and rewrote the file anyway. All 46
    parts came back byte-identical in CONTENT, but every zip entry was reordered
    and every timestamp reset — and a full python-pptx repackage is exactly where
    the parts it does not model get dropped. Nothing was gained by writing.

    Both no-change shapes are checked: nothing to find at all, and something found
    but out of reach. The second one must ALSO keep printing its `[!]` line — an
    early return that skips the write is one line away from skipping the warning
    that made this run worth reporting.
    """
    out = []
    r = ctx["replace_noop"]
    if not ctx["noop_untouched"]:
        out.append("a 0-replacement run rewrote the input file")
    if r["exit"] != 0:
        out.append(f"exit {r['exit']} on a run with nothing to replace")
    if not re.search(r"(?m)^Unchanged ", r["out"]):
        out.append("it did not say the file was left alone")
    if re.search(r"(?m)^Saved ", r["out"]):
        out.append("it still claims it saved something")
    if reported_count(r) != 0:
        out.append(f"reported {reported_count(r)} replacements, expected 0")
    u = ctx["replace_unreachable"]
    if not ctx["unreachable_untouched"]:
        out.append("a run whose only match was out of reach rewrote the input file")
    if not re.search(r"(?m)^Unchanged ", u["out"]):
        out.append("the out-of-reach case did not say the file was left alone")
    if reported_missed(u["out"], TABLE_ONLY) != 1:
        out.append("skipping the write also swallowed the [!] out-of-reach line")
    return out


@check("W2", "an in-place run that DID change something writes the file")
def w2(ctx: dict) -> list[str]:
    """The other direction of N4, and the reason N4 cannot be satisfied by simply
    never writing in place. Every other assertion in this gate passes --out, so
    without this one the default path is only ever asserted to do nothing."""
    out = []
    if not ctx["inplace_written"]:
        out.append("an in-place replacement did not write the file")
    if ctx["inplace_applied"] != 1:
        out.append(f"the in-place replacement did not apply ({ctx['inplace_applied']})")
    if not re.search(r"(?m)^Saved ", ctx["replace_inplace"]["out"]):
        out.append("it did not report saving the file")
    return out


@check("W1", "it does replace what it CAN see")
def w1(ctx: dict) -> list[str]:
    """Without this, every limit assertion above is also satisfied by a tool that
    replaces nothing at all — and a tool that finds nothing gets reported as broken
    within a minute, while one that finds some things gets trusted forever."""
    out = []
    if ctx["single_applied"] != 1:
        out.append(f"a single-run phrase was not replaced ({ctx['single_applied']})")
    if reported_count(ctx["replace_single"]) != 1:
        out.append(f"reported {reported_count(ctx['replace_single'])}, expected 1")
    return out


@check("X1", "reading a .pptx with no slideLayout survives with one sentence")
def x1(ctx: dict) -> list[str]:
    r = ctx["read_nolayout"]
    out = []
    if "Traceback" in r["err"]:
        out.append("bare traceback on a layout-less .pptx")
    if r["exit"] != 0:
        out.append(f"exit {r['exit']} — the text is still readable, "
                   f"a decorative layout name must not cost the whole document")
    if "(no layout)" not in r["out"]:
        out.append("the missing layout is not reported as such")
    if PROBE not in r["out"]:
        out.append("the slide text was lost along with the layout")
    return out


@check("X2", "adding a slide to a master-less .pptx fails with one sentence")
def x2(ctx: dict) -> list[str]:
    r = ctx["add_nolayout"]
    out = []
    if "Traceback" in r["err"]:
        out.append("bare traceback — the bounds check itself crashed")
    if r["exit"] == 0:
        out.append("reported success on a file that has no slide master")
    if "no slide master" not in r["err"]:
        out.append("the reason is not stated")
    return out


@check("X3", "adversarial input never produces a traceback")
def x3(ctx: dict) -> list[str]:
    out = []
    for label, res in ctx["adversarial"].items():
        for name, r in res.items():
            if "Traceback" in r["err"]:
                out.append(f"{name} on 「{label}」: bare traceback")
            if r["exit"] == 0:
                out.append(f"{name} on 「{label}」: exit 0")
    return out


@check("A1", "--layout out of range is refused and the valid range is named")
def a1(ctx: dict) -> list[str]:
    out = []
    for label, r in (("99", ctx["layout_high"]), ("-1", ctx["layout_neg"])):
        if r["exit"] == 0:
            out.append(f"--layout {label} was accepted")
        if "have 0.." not in r["err"]:
            out.append(f"--layout {label}: the valid range is not named")
    return out


@check("C1", "Chinese output survives a captured stdout on an ANSI code page")
def c1(ctx: dict) -> list[str]:
    """A Windows-only product defect, found by this gate's FIRST CI run.

    Windows encodes a captured stdout in the machine's ANSI code page, and Python
    only defaults to UTF-8 from 3.15 (PEP 686); CI pins 3.11. Without the two-line
    reconfigure at the top of each script, `pptx_read.py` exits 1 with
    UnicodeEncodeError on the first Chinese character it prints — and an agent
    ALWAYS captures stdout, so on Windows the skill simply did not work on any
    deck with Chinese in it. It shipped that way from the doc-edit days: no gate
    had ever covered these two scripts.
    """
    out = []
    r = ctx["read_ansi"]
    if r["exit"] != 0:
        out.append(f"pptx_read exited {r['exit']} under an ANSI code page")
    if "UnicodeEncodeError" in r["err"]:
        out.append("pptx_read: UnicodeEncodeError on captured stdout")
    if PROBE not in r["out"]:
        out.append("pptx_read: the Chinese slide text did not survive")
    e = ctx["replace_ansi"]
    if e["exit"] != 0:
        out.append(f"pptx_edit exited {e['exit']} under an ANSI code page")
    if "UnicodeEncodeError" in e["err"]:
        out.append("pptx_edit: UnicodeEncodeError on captured stdout")
    return out


@check("E1", "--out leaves the input byte-identical")
def e1(ctx: dict) -> list[str]:
    return [] if ctx["input_untouched"] else ["--out still modified the input file"]


# ── negative controls ─────────────────────────────────────────────────────────
def patched(work: Path, old: str, new: str, name: str) -> Path:
    """A copy of the skill's scripts with one edit applied.

    Raises if the anchor is not found exactly once — a control arm that silently
    fails to apply looks identical to one that applied and did nothing.
    """
    dest = work / f"patched-{name}"
    if dest.exists():
        shutil.rmtree(dest)
    shutil.copytree(SCRIPTS, dest)
    hits = 0
    for py in dest.glob("*.py"):
        text = py.read_text(encoding="utf-8")
        if old in text:
            hits += text.count(old)
            py.write_text(text.replace(old, new), encoding="utf-8")
    if hits != 1:
        raise SystemExit(f"control {name!r}: anchor matched {hits} times, expected 1 "
                         f"— the control did not replicate the defect")
    return dest


def patched_in(script_dir: Path, old: str, new: str) -> int:
    """Apply one more edit to an ALREADY patched copy. Returns the hit count so the
    caller can refuse a control that did not fully apply."""
    hits = 0
    for py in script_dir.glob("*.py"):
        text = py.read_text(encoding="utf-8")
        if old in text:
            hits += text.count(old)
            py.write_text(text.replace(old, new), encoding="utf-8")
    return hits


def live_read_crashes(work: Path):
    """Restore the implementation that shipped until 2026-08-04, verbatim."""
    d = patched(work, """        try:
            layout = slide.slide_layout.name
        except Exception:  # noqa: BLE001
            layout = "(no layout)"
        slides.append({"index": idx, "layout": layout, "texts": texts,
                       "unread": _census(slide)})""",
        """        slides.append({"index": idx, "layout": slide.slide_layout.name, "texts": texts,
                       "unread": _census(slide)})""",
        "read-crash")
    return collect(work / "c1", d)


def live_bounds_crashes(work: Path):
    d = patched(work, """        try:
            n_layouts = len(prs.slide_layouts)
        except Exception as exc:  # noqa: BLE001
            print(f"Cannot add a slide to {args.file}: it has no slide master ({exc})", file=sys.stderr)
            return 1
        if args.layout < 0 or args.layout >= n_layouts:
            print(f"Bad --layout {args.layout} (have 0..{n_layouts - 1})", file=sys.stderr)
            return 1""",
        """        if args.layout < 0 or args.layout >= len(prs.slide_layouts):
            print(f"Bad --layout {args.layout} (have 0..{len(prs.slide_layouts) - 1})", file=sys.stderr)
            return 1""",
        "bounds-crash")
    return collect(work / "c2", d)


def live_replaces_nothing(work: Path):
    # Anchored on the ASSIGNMENT, not on the bare `if old in run.text:` — the same
    # condition now also appears (more deeply indented) in the reach census, and a
    # short anchor is a SUBSTRING of the longer line. Matching both would have
    # disabled the census too, turning a clean control into a cascade.
    d = patched(work, """            if old in run.text:
                run.text = run.text.replace(old, new)
                n += 1""",
        """            if False:
                run.text = run.text.replace(old, new)
                n += 1""", "no-op")
    return collect(work / "c3", d)


def live_out_ignored(work: Path):
    d = patched(work, "    out = args.out or args.file", "    out = args.file", "out-ignored")
    return collect(work / "c4", d)


def live_read_walks_tables(work: Path):
    d = patched(work,
        "        texts = [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()]",
        "        texts = [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()]\n"
        "        for _sh in slide.shapes:\n"
        "            if getattr(_sh, 'has_table', False) and _sh.has_table:\n"
        "                texts += [c.text for r in _sh.table.rows for c in r.cells if c.text.strip()]",
        "walks-tables")
    return collect(work / "c5", d)


def live_no_utf8_guard(work: Path):
    """Remove the reconfigure block from pptx_read — i.e. the state the skill
    shipped in until 2026-08-04. On a UTF-8 machine this changes nothing at all,
    which is the whole point: only the forced code page tells the two apart."""
    d = patched(work, '''for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        try:
            _stream.reconfigure(encoding="utf-8")
        except (ValueError, OSError):        # already detached / not reconfigurable
            pass

def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def _unread_kind(shape):''',
        '''def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def _unread_kind(shape):''',
        "no-utf8-guard")
    return collect(work / "c9", d)


def live_read_silent_on_unread(work: Path):
    """pptx_read as it shipped until 2026-08-16: it skipped tables and groups and
    said nothing, so a table-only slide printed exactly like a title-only one."""
    # Anchored on BOTH halves of the reporting — per-slide marker and the closing
    # summary. A first cut disabled only the per-slide half, and the summary line
    # went on printing the correct counts, so N1 stayed green against a control
    # that was supposed to break it. The control was wrong, and it also exposed a
    # hole in N1 (see the line-anchored regex there).
    d = patched(work, """            u = s["unread"]
            if u["table"] or u["group"]:""",
        """            u = {"table": 0, "group": 0}
            if False:""",
        "read-silent")
    quiet = patched_in(d, """        if totals["table"] or totals["group"]:""",
                       """        if False:""")
    if not quiet:
        raise SystemExit("control 'read-silent': the summary half did not apply")
    return collect(work / "c10", d)


def live_edit_silent_on_missed(work: Path):
    """pptx_edit as it shipped until 2026-08-16: `replacements: N` and not one word
    about what it could see but could not reach."""
    d = patched(work, "    for old, n in missed:", "    for old, n in []:", "edit-silent")
    return collect(work / "c11", d)


def live_always_rewrites(work: Path):
    """pptx_edit as it shipped until 2026-08-16 evening: it saved unconditionally,
    so asking for a replacement that matched nothing still repackaged the file."""
    d = patched(work, "    if args.out or total or added:", "    if True:",
                "always-rewrites")
    return collect(work / "c13", d)


def live_never_writes_in_place(work: Path):
    """The overcorrection: skip the write whenever there is no --out, changes or
    not. N4 cannot tell this apart from the fix — W2 is the only thing that can,
    which is the whole reason it exists."""
    d = patched(work, "    if args.out or total or added:", "    if args.out:",
                "never-writes-in-place")
    return collect(work / "c14", d)


def live_census_skips_groups(work: Path):
    """A half-blind census: it walks tables and forgets groups.

    This is the likeliest way to get the fix wrong, and the table-only fixture
    cannot see it — which is why N1/N2 each carry a group case.
    """
    d = patched(work, """        if not shape.has_text_frame and getattr(shape, "shapes", None) is not None:
            return "group\"""",
        """        if False:
            return "group\"""",
        "census-skips-groups")
    return collect(work / "c12", d)


def fixture_no_table(work: Path):
    return collect(work / "c6", SCRIPTS, with_table=False)


def fixture_one_run(work: Path):
    return collect(work / "c7", SCRIPTS, split_phrase=False)


def fixture_has_layouts(work: Path):
    return collect(work / "c8", SCRIPTS, layoutless_has_layouts=True)


FLAWS = [
    ("LIVE: pptx_read as it shipped until 2026-08-04 (unguarded slide_layout)",
     live_read_crashes, {"X1"}, ""),
    ("LIVE: pptx_edit as it shipped until 2026-08-04 (bounds check reaches the master)",
     live_bounds_crashes, {"X2"}, ""),
    ("LIVE: --replace matches nothing at all", live_replaces_nothing,
     {"W1", "L4", "W2"},
     "L4 also fires, and must: 'the text-frame copy was replaced' is the half of "
     "L4 that a do-nothing implementation breaks. W2 also fires, and must: with "
     "every replacement disabled the in-place run has nothing to write, so the "
     "file is correctly left alone — for the wrong reason. That is the second "
     "detector of a do-nothing tool, not a cascade to be suppressed"),
    ("LIVE: --out is ignored and it writes in place", live_out_ignored,
     {"E1", "L4", "W1"},
     "L4/W1 also fire, and must: nothing is ever written to the --out path, so "
     "every assertion that reads that artifact has no artifact to read. Declared "
     "rather than explained away — a cascade note does not suppress anything here, "
     "an undeclared check that fires is still a failure"),
    ("LIVE: pptx_read learns to walk tables", live_read_walks_tables, {"L1"},
     "L4 does NOT fire: reading tables changes nothing about what --replace does"),
    ("LIVE: pptx_read without the UTF-8 reconfigure (as it shipped until 2026-08-04)",
     live_no_utf8_guard,
     {"C1"} if HOST_CAPTURES_UTF8 else {"C1", "V0", "X1"},
     f"this host hands a captured child stdout `{CAPTURED_ENC}`. "
     + ("On a UTF-8 host V0/X1 do NOT fire — the guarded and unguarded scripts behave "
        "identically and only the forced ANSI code page in C1 can tell them apart."
        if HOST_CAPTURES_UTF8 else
        "On a host whose default captured encoding is ALREADY the hostile code page, "
        "removing the guard breaks EVERY run of the script, so V0/X1 fire too — that "
        "is the real Windows defect, not a cascade to be explained away.")),
    ("LIVE: pptx_read as it shipped until 2026-08-16 (silent about unread containers)",
     live_read_silent_on_unread, {"N1"},
     "L1 does NOT fire: staying silent about tables is not the same as printing "
     "their text, and L1 only watches the text"),
    ("LIVE: pptx_edit as it shipped until 2026-08-16 (silent about unreachable matches)",
     live_edit_silent_on_missed, {"N2", "N4"},
     "L4 does NOT fire, and must not: the silent version replaced exactly the same "
     "runs. L4 pins the reach, N2 pins the admission — this control is what proves "
     "they are two different properties. N4 also fires, by design and declared: one "
     "of its sub-assertions says the skipped write must not swallow the [!] line, "
     "and this control removes that line. The overlap is the guard against a future "
     "early return in the no-change branch taking the warning with it"),
    ("LIVE: pptx_edit as it shipped until 2026-08-16 evening (rewrites even with "
     "nothing to change)", live_always_rewrites, {"N4"},
     "W2 does NOT fire: writing too often still writes when there IS a change. "
     "N2 does not fire either — the unconditional save never touched the [!] line"),
    ("LIVE: the in-place write is skipped even when something DID change",
     live_never_writes_in_place, {"W2"},
     "N4 does NOT fire, and that is the point: a tool that never writes in place "
     "satisfies every no-change assertion. Only the other direction sees it"),
    ("LIVE: the census walks tables but forgets groups", live_census_skips_groups,
     {"N1"},
     "N2 does NOT fire: pptx_edit's own reach census is separate code, so a blind "
     "spot in pptx_read's does not travel to it"),
    ("CONTROL: fixture loses the table", fixture_no_table, {"V0", "L4", "N4"},
     "L4 also fires: with no second copy there is nothing for it to count. "
     "N1/N2 do NOT fire: both compare the report against the fixture's real "
     "contents, and a fixture with no table is honestly reported as having none. "
     "N4 fires on its out-of-reach half only — with no table there is no "
     "out-of-reach match to warn about, i.e. that half of N4 has lost its subject, "
     "which is exactly what this fixture control is for"),
    ("CONTROL: fixture phrase is no longer split across runs", fixture_one_run,
     {"V0", "L3"}, "L3 also fires: a single-run phrase DOES match, which is the "
                   "very thing L3 asserts cannot happen"),
    ("CONTROL: the layout-less fixture actually has layouts", fixture_has_layouts,
     {"V0", "X1", "X2"},
     "X1 and X2 both fire, and must: with a real master present there is no "
     "'(no layout)' for X1 to find, and --add-slide legitimately SUCCEEDS, which "
     "is exactly the 'reported success' branch X2 watches"),
]


def fired(ctx: dict) -> dict[str, list[str]]:
    return {cid: f for cid, c in CHECKS.items() if (f := c["fn"](ctx))}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    try:
        import pptx  # noqa: F401
    except ImportError:
        print("python-pptx is not installed — this gate cannot run "
              "(it is a required dependency of the skill under test)", file=sys.stderr)
        return 1

    results, passed, failed = [], 0, 0
    with tempfile.TemporaryDirectory(prefix="pptx-edit-gate-") as td:
        work = Path(td)
        (work / "base").mkdir()
        base = collect(work / "base", SCRIPTS)

        clean = fired(base)
        ok = not clean
        results.append({"case": "real output of the real scripts", "expect": "silence",
                        "ok": ok, "detail": [f"{k}: {v[0]}" for k, v in clean.items()]})
        passed, failed = (1, 0) if ok else (0, 1)

        for name, mutate, expected, cascade in FLAWS:
            sub = work / f"flaw-{len(results)}"
            sub.mkdir()
            ctx = mutate(sub)
            got = fired(ctx)
            unexpected, missing = set(got) - expected, expected - set(got)
            ok = not unexpected and not missing
            detail = []
            if missing:
                detail.append(f"did NOT fire: {sorted(missing)}")
            if unexpected:
                detail.append(f"unexpectedly fired: {sorted(unexpected)}")
            results.append({"case": name, "expect": sorted(expected),
                            "fired": sorted(got), "ok": ok, "detail": detail,
                            "cascade_note": cascade})
            passed, failed = (passed + 1, failed) if ok else (passed, failed + 1)

    if args.json:
        print(json.dumps({"results": results, "assertions": len(CHECKS),
                          "passed": passed, "failed": failed}, ensure_ascii=False, indent=2))
    else:
        for r in results:
            print(f"{'PASS' if r['ok'] else 'FAIL'}  {r['case']}")
            if r.get("cascade_note"):
                print(f"        note: {r['cascade_note']}")
            for d in r["detail"]:
                print(f"        {d}")
        print(f"\n[pptx-edit-skill] {passed} passed, {failed} failed, "
              f"{len(CHECKS)} assertions "
              f"(captured-stdout encoding on this host: {CAPTURED_ENC})")
    return 0 if failed == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
