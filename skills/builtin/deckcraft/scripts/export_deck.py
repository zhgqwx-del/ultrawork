#!/usr/bin/env python3
"""Export a built deck.html to PDF / per-page PNGs / image-type PPTX via headless Chrome/Edge.

    python3 export_deck.py <project_dir> [--pdf] [--shots] [--pptx] [--out <dir>]

  --pdf    write <out>/deck.pdf  (default action if no flag given)
  --shots  write <out>/shots/pNN.png per page (used by the visual-QA step)
  --pptx   write <out>/deck.pptx — image-type: each slide is a full-bleed
           screenshot (text NOT editable in PowerPoint; state this when
           delivering). Implies --shots. Needs the python-pptx pip library.
  --out    output dir, default <project_dir>/export
  --publish <dir>  copy the final deliverables (deck.html + produced pdf/pptx)
           into <dir> named <project-name>.html/.pdf/.pptx — the workspace-visible
           delivery step. The project dir itself should be a dot-directory
           (.deckcraft/<name>) so intermediates stay out of the artifacts panel.

Requires deck.html built by build_deck.py. Stdlib + a Chromium browser
(located by find_chrome.py); python-pptx only for --pptx. Cross-platform:
no shell, no hardcoded /tmp.
"""
from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from find_chrome import find_browser
from find_node import find_node

from console_encoding import configure_utf8_stdio

configure_utf8_stdio()

CHROME_BASE = ["--headless=new", "--disable-gpu", "--hide-scrollbars", "--no-first-run", "--no-default-browser-check"]
TIMEOUT_S = 120

NODE_MISSING_MSG = (
    "ERROR: --pptx-editable needs a Node.js runtime and none was found.\n"
    "  The editable-pptx export uses the app's embedded Node — install it via\n"
    "  Settings → Browser dependencies (it downloads a private Node runtime).\n"
    "  Or export an image-type pptx instead: export_deck.py <project> --pptx"
)


def run_chrome(browser: str, args: list[str]) -> None:
    kwargs = {}
    if sys.platform.startswith("win"):
        # GUI-subsystem parents must not flash console windows (app convention, ADR-054)
        kwargs["creationflags"] = 0x08000000  # CREATE_NO_WINDOW
    try:
        subprocess.run([browser, *CHROME_BASE, *args], check=True, timeout=TIMEOUT_S,
                       stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, **kwargs)
    except subprocess.CalledProcessError as e:
        raise SystemExit(f"ERROR: browser exited {e.returncode} for: {' '.join(args[:2])}")
    except subprocess.TimeoutExpired:
        raise SystemExit(f"ERROR: browser timed out after {TIMEOUT_S}s for: {' '.join(args[:2])}")


def split_pages(deck_html: str) -> tuple[str, list[str]]:
    """Return (head_html, [section_html...]) — head is everything before the first slide.

    HTML comments are stripped first so documentation comments can never
    contribute phantom sections or corrupt the head slice. Sections are matched
    on the bare <section tag (not an exact attribute order) — validate_deck
    guarantees one non-nested <section per page, so any attribute layout the
    model produced splits correctly.
    """
    deck_html = re.sub(r"<!--.*?-->", "", deck_html, flags=re.S)
    sections = re.findall(r"<section\b.*?</section>", deck_html, re.S)
    head = re.split(r"<section\b", deck_html, 1)[0]
    # head still contains the opening <div class="stage">; close it per single page
    return head, sections


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("project_dir")
    ap.add_argument("--pdf", action="store_true")
    ap.add_argument("--shots", action="store_true")
    ap.add_argument("--pptx", action="store_true")
    ap.add_argument("--pptx-editable", dest="pptx_editable", action="store_true",
                    help="editable .pptx (P2b): DOM elements → text boxes/shapes "
                         "via Node+pptxgenjs; needs a Node runtime")
    ap.add_argument("--out", default=None)
    ap.add_argument("--publish", default=None)
    a = ap.parse_args()
    if a.pptx:
        a.shots = True  # image pptx is assembled from the per-page screenshots

    proj = Path(a.project_dir)
    deck = proj / "deck.html"
    if not deck.is_file():
        print(f"ERROR: {deck} not found — run build_deck.py first", file=sys.stderr)
        return 1
    browser = find_browser()
    if not browser:
        print("ERROR: no Chrome/Edge/Chromium found. Install Google Chrome or Microsoft Edge.", file=sys.stderr)
        return 1

    out_dir = Path(a.out) if a.out else proj / "export"
    out_dir.mkdir(parents=True, exist_ok=True)
    # explicit --pdf, or no action flag at all → default to PDF
    do_pdf = a.pdf or not (a.shots or a.pptx_editable)

    if a.pptx:  # fail fast BEFORE spending ~1s/page on screenshots
        try:
            import pptx  # noqa: F401
        except ImportError:
            print("ERROR: --pptx needs the python-pptx library (pip install python-pptx)",
                  file=sys.stderr)
            return 1

    node_bin = None
    if a.pptx_editable:  # fail fast on missing Node BEFORE any Chrome work
        node_bin = find_node()
        if not node_bin:
            print(NODE_MISSING_MSG, file=sys.stderr)
            return 1

    if do_pdf:
        pdf = out_dir / "deck.pdf"
        run_chrome(browser, [f"--print-to-pdf={pdf}", "--no-pdf-header-footer", deck.resolve().as_uri()])
        print(f"OK: {pdf}")

    n_pages = 0
    if a.shots:
        shots = out_dir / "shots"
        # stale shots from a previous run (e.g. after dropping a page) would leak
        # into the pptx and the visual QA — always start from an empty dir
        if shots.is_dir():
            shutil.rmtree(shots)
        shots.mkdir(parents=True)
        html = deck.read_text(encoding="utf-8")
        head, sections = split_pages(html)
        if not sections:
            print("ERROR: no slides found in deck.html", file=sys.stderr)
            return 1
        # flatten margins via a style override (attribute-order independent) so
        # each page renders at exactly 1280x720
        flat = "<style>.stage{margin:0}.slide{margin:0}</style>"
        # 2x device pixels for the pptx path — projector/zoom sharpness
        scale = ["--force-device-scale-factor=2"] if a.pptx else []
        with tempfile.TemporaryDirectory() as td:
            for i, sec in enumerate(sections, 1):
                page = Path(td) / f"p{i:02d}.html"
                page.write_text(head + flat + sec + "</div></body></html>", encoding="utf-8")
                run_chrome(browser, ["--window-size=1280,720", *scale,
                                     f"--screenshot={shots / f'p{i:02d}.png'}", page.resolve().as_uri()])
        n_pages = len(sections)
        print(f"OK: {shots} ({n_pages} pages)")

    if a.pptx:
        pptx_out = build_image_pptx(out_dir / "shots", out_dir / "deck.pptx",
                                    speaker_notes(proj), expected=n_pages)
        print(f"OK: {pptx_out} (image-type — text not editable in PowerPoint)")

    if a.pptx_editable:
        edit_out, raster_report, text_lost = build_editable_pptx(proj, out_dir, browser, node_bin)
        n_raster = sum(raster_report.values())
        print(f"OK: {edit_out} (editable — text/shapes are PowerPoint-native)")
        if n_raster:
            pages = ", ".join(f"第 {i} 页 {c} 个" for i, c in sorted(raster_report.items()))
            extra = f"，其中约 {text_lost} 处原本可编辑的文字一并被栅格化" if text_lost else ""
            print(f"NOTE: {n_raster} 个元素无法翻译为可编辑对象，已栅格化为图片（{pages}）{extra}"
                  f"——这些元素在 PowerPoint 里不可编辑。")

    if a.publish:
        dest = Path(a.publish)
        dest.mkdir(parents=True, exist_ok=True)
        # visible name = project dir name with any leading dots stripped, so a
        # published deliverable is never itself a hidden dotfile
        name = proj.resolve().name.lstrip(".") or "deck"
        # only deliverables THIS run produced — a stale deck.pptx from an earlier
        # experiment must not ship when the user chose HTML+PDF this time
        candidates = [(deck, ".html")]
        if do_pdf:
            candidates.append((out_dir / "deck.pdf", ".pdf"))
        # editable pptx claims the plain <name>.pptx; if the user also built the
        # image type this run, it publishes under <name>-image.pptx so neither is lost
        if a.pptx_editable:
            candidates.append((out_dir / "deck-editable.pptx", ".pptx"))
            if a.pptx:
                candidates.append((out_dir / "deck.pptx", "-image.pptx"))
        elif a.pptx:
            candidates.append((out_dir / "deck.pptx", ".pptx"))
        published = []
        for src, ext in candidates:
            if src.is_file():
                target = dest / f"{name}{ext}"
                if target.exists():
                    print(f"WARN: overwriting existing {target}")
                shutil.copyfile(src, target)
                published.append(str(target))
        for t in published:
            print(f"PUBLISHED: {t}")
    return 0


def _shoot_pages_2x(browser: str, head: str, sections: list[str],
                    indices: list[int], td: Path) -> dict[int, Path]:
    """Screenshot the given 1-based pages at 2x device pixels (2560x1440) into
    `td`, returning {index: png_path}. Used only to crop rasterized regions —
    same per-page isolation as the main shots path (drops the screen-fit script)."""
    flat = "<style>.stage{margin:0}.slide{margin:0}</style>"
    out: dict[int, Path] = {}
    for i in indices:
        page = td / f"crop-p{i:02d}.html"
        page.write_text(head + flat + sections[i - 1] + "</div></body></html>",
                        encoding="utf-8")
        png = td / f"crop-p{i:02d}.png"
        run_chrome(browser, ["--window-size=1280,720", "--force-device-scale-factor=2",
                             f"--screenshot={png}", page.resolve().as_uri()])
        out[i] = png
    return out


def build_editable_pptx(proj: Path, out_dir: Path, browser: str,
                        node_bin: str) -> tuple[Path, dict[int, int], int]:
    """Extract deck.html into positioned primitives and assemble an editable
    .pptx via Node+pptxgenjs. Elements that cannot be translated faithfully
    (inline SVG, remote/relative images, background-image) are rasterized from a
    2x screenshot and reported. Returns (pptx_path, {page_index: rasterized_count},
    total_text_boxes_lost_to_rasterization).

    All intermediates (layout.json, raster/) live under out_dir (the export/
    dir, which pack/git already exclude) so they never leak into the artifacts
    panel or the committed example baseline."""
    # lazy import avoids the extract_layout ↔ export_deck circular import at module load
    from extract_layout import extract

    layout = extract(proj, browser)
    # carry speaker notes so the editable pptx keeps presenter-view parity with the
    # image-type pptx (both read outline.json); empty when there's no outline
    notes = speaker_notes(proj)
    for pg in layout["pages"]:
        note = notes.get(pg["index"], "")
        if note:
            pg["notes"] = note
    rasters = [(pg, el) for pg in layout["pages"]
               for el in pg["elements"] if el["type"] == "raster"]
    raster_report: dict[int, int] = {}
    text_lost = 0

    if rasters:
        try:
            from PIL import Image
        except ImportError:
            raise SystemExit(
                "ERROR: rasterizing non-editable elements needs Pillow "
                "(pip install Pillow — it also ships with python-pptx)")
        raster_dir = out_dir / "raster"
        if raster_dir.is_dir():
            shutil.rmtree(raster_dir)
        raster_dir.mkdir(parents=True)
        head, sections = split_pages((proj / "deck.html").read_text(encoding="utf-8"))
        pages_with_raster = sorted({pg["index"] for pg, _ in rasters})
        # group by page so each 2x screenshot is opened exactly once, not per element
        by_page: dict[int, list] = {}
        for pg, el in rasters:
            by_page.setdefault(pg["index"], []).append((pg, el))
        counters: dict[int, int] = {}
        dropped: list = []  # off-canvas elements — removed so none is left as raster
        with tempfile.TemporaryDirectory() as td:
            shots = _shoot_pages_2x(browser, head, sections, pages_with_raster, Path(td))
            for idx, entries in by_page.items():
                with Image.open(shots[idx]) as img:
                    w, h = img.width, img.height
                    for pg, el in entries:
                        b = el["box"]
                        left, top = max(0, int(b["x"] * 2)), max(0, int(b["y"] * 2))
                        right, bottom = min(w, int((b["x"] + b["w"]) * 2)), min(h, int((b["y"] + b["h"]) * 2))
                        if right <= left or bottom <= top:
                            # (partly) outside the 1280x720 canvas — only reachable if
                            # the overflow gate was skipped; drop it so it is neither a
                            # crash (inverted crop box) nor a fatal "unresolved raster"
                            dropped.append((pg, el))
                            continue
                        k = counters.get(idx, 0)
                        counters[idx] = k + 1
                        png = raster_dir / f"p{idx:02d}-e{k}.png"
                        img.crop((left, top, right, bottom)).save(png)
                        text_lost += int(el.get("textLost") or 0)
                        # rewrite the raster element in-place into an image the assembler places
                        el.clear()
                        el.update({"type": "image", "box": b,
                                   "src": str(png.resolve()), "rasterized": True})
                        raster_report[idx] = raster_report.get(idx, 0) + 1
        for pg, el in dropped:
            pg["elements"].remove(el)

    layout_path = out_dir / "layout.json"
    layout_path.write_text(json.dumps(layout, ensure_ascii=False), encoding="utf-8")

    out_pptx = out_dir / "deck-editable.pptx"
    assemble = Path(__file__).resolve().parent / "html2pptx" / "assemble.mjs"
    kwargs = {}
    if sys.platform.startswith("win"):
        kwargs["creationflags"] = 0x08000000  # CREATE_NO_WINDOW (ADR-054)
    r = subprocess.run([node_bin, str(assemble), str(layout_path), str(out_pptx)],
                       capture_output=True, text=True, encoding="utf-8",
                       errors="replace", timeout=180, **kwargs)
    if r.returncode != 0:
        raise SystemExit(f"ERROR: pptx assembler failed (exit {r.returncode}):\n"
                         f"{r.stdout}{r.stderr}")
    return out_pptx, raster_report, text_lost


def speaker_notes(proj: Path) -> dict[int, str]:
    """Per-page speaker notes from outline.json (empty dict when absent)."""
    outline_f = proj / "outline.json"
    if not outline_f.is_file():
        return {}
    try:
        data = json.loads(outline_f.read_text(encoding="utf-8"))
        return {s["index"]: str(s.get("speaker_notes") or "")
                for s in data.get("slides", []) if "index" in s}
    except (json.JSONDecodeError, TypeError, KeyError):
        return {}


def build_image_pptx(shots_dir: Path, out_path: Path, notes: dict[int, str],
                     expected: int = 0) -> Path:
    """Assemble a 16:9 image-type .pptx: one full-bleed screenshot per slide,
    speaker notes attached — presenter view stays fully usable even though
    the slide surface itself is an image."""
    from pptx import Presentation
    from pptx.util import Inches

    def shot_no(p: Path) -> int:
        m = re.search(r"(\d+)", p.stem)
        return int(m.group(1)) if m else 0

    # numeric sort — lexicographic would misplace p100 before p11
    shots = sorted(shots_dir.glob("p*.png"), key=shot_no)
    if not shots:
        raise SystemExit(f"ERROR: no screenshots in {shots_dir}")
    if expected and len(shots) != expected:
        raise SystemExit(f"ERROR: {len(shots)} screenshots != {expected} deck pages "
                         f"in {shots_dir} — stale files?")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout — no placeholders
    for shot in shots:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(shot), 0, 0,
                                 width=prs.slide_width, height=prs.slide_height)
        # note alignment relies on validate_deck E5 enforcing contiguous 1..N page
        # numbers == outline indices; the gate runs before export in the workflow
        note = notes.get(shot_no(shot), "")
        if note:
            slide.notes_slide.notes_text_frame.text = note
    prs.save(str(out_path))
    return out_path


if __name__ == "__main__":
    sys.exit(main())
