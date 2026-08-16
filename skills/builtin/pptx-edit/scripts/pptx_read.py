#!/usr/bin/env python3
"""Read a .pptx: dump slide outline (text of each shape). ultrawork pptx-edit skill.

Only top-level shapes with a text frame are walked: table cells (GraphicFrame) and
grouped shapes (GroupShape) are NOT reported. Measured, and stated in SKILL.md
「限制」 — do not let a caller assume an empty slide means an empty slide.

What this script does NOT do is stay quiet about it. It reports that those
containers EXIST (never their text — that stays out of scope, and the gate asserts
it in both directions), because without that line a slide whose only content is a
table is byte-for-byte indistinguishable from a slide that carries nothing but a
title. Measured in L4 review §四: a model recovered the table on 季度汇报.pptx only
because that slide's TITLE happened to read 「分部收入（表格）」 — the fixture leaked
the answer. Rename the slide and the same run loses the whole table silently.
"""
from __future__ import annotations
import argparse, json, sys

# ── stdout must be UTF-8 on every platform, and on Windows it is not ─────────
# This script prints whatever text the slides carry, which for this project's users
# is Chinese. Windows encodes a CAPTURED stdout in the machine's ANSI code page
# (cp1252 on a western install), and Python only defaults to UTF-8 from 3.15
# (PEP 686) — CI pins 3.11. Measured: WITHOUT the lines below this script exits 1
# with `UnicodeEncodeError: 'charmap' codec can't encode character '\u7b2c'` the
# moment its output is captured — which is exactly how an agent calls it. The other
# office skills each carry this in a shared module; this one has no shared module,
# so both entry points carry it. Found by the first CI run of this skill's gate,
# not by any local run: it cannot be seen on macOS or Linux.
for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        try:
            _stream.reconfigure(encoding="utf-8")
        except (ValueError, OSError):        # already detached / not reconfigurable
            pass

def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def _unread_kind(shape):
    """Which unreadable text container this top-level shape is, if any.

    A table is a GraphicFrame and a group is a GroupShape; both answer False to
    `has_text_frame`, which is why the walk below skips them. A picture also has no
    text frame but carries no text either, so it is not counted — the number has to
    mean "text you are not seeing", or it trains the reader to ignore it.

    Every probe is guarded: python-pptx raises from several of these properties on
    shapes built by non-PowerPoint generators, and a census must never be the thing
    that turns a readable deck into a traceback.
    """
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            return "table"
    except Exception:  # noqa: BLE001
        return None
    try:
        if not shape.has_text_frame and getattr(shape, "shapes", None) is not None:
            return "group"
    except Exception:  # noqa: BLE001
        return None
    return None


def _census(slide):
    n = {"table": 0, "group": 0}
    try:
        shapes = list(slide.shapes)
    except Exception:  # noqa: BLE001
        return n
    for sh in shapes:
        kind = _unread_kind(sh)
        if kind:
            n[kind] += 1
    return n


def main(argv):
    ap = argparse.ArgumentParser(description="Read slide outline from a .pptx")
    ap.add_argument("file")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args(argv)
    pptx = _require()
    try:
        prs = pptx.Presentation(args.file)
    except Exception as exc:  # noqa: BLE001
        print(f"Error opening {args.file}: {exc}", file=sys.stderr)
        return 1
    slides = []
    for idx, slide in enumerate(prs.slides):
        texts = [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()]
        # A slide is not required to carry a slideLayout relationship, and files
        # written by minimal OOXML generators (rather than PowerPoint) often don't:
        # python-pptx then raises KeyError deep inside the relationship lookup.
        # The layout NAME is decoration here — the text is the payload — so a
        # missing one must not cost the caller the whole document (measured on a
        # real 10-part .pptx that had no ppt/slideLayouts at all).
        try:
            layout = slide.slide_layout.name
        except Exception:  # noqa: BLE001
            layout = "(no layout)"
        slides.append({"index": idx, "layout": layout, "texts": texts,
                       "unread": _census(slide)})
    totals = {k: sum(s["unread"][k] for s in slides) for k in ("table", "group")}
    if args.json:
        print(json.dumps({"slides": slides, "unread_total": totals},
                         ensure_ascii=False, indent=2))
    else:
        for s in slides:
            print(f"=== slide {s['index']} ({s['layout']}) ===")
            for t in s["texts"]:
                print(t)
            u = s["unread"]
            if u["table"] or u["group"]:
                bits = ", ".join(f"{label} ×{u[k]}" for k, label in
                                 (("table", "表格"), ("group", "组合")) if u[k])
                print(f"[unread] 本页另有本脚本读不到的元素：{bits}")
        if totals["table"] or totals["group"]:
            # One concrete next action, not just a status field. Measured in L4
            # review §二: fields that only state a condition were dropped by the
            # model three times out of three, while the one note naming an action
            # to run was acted on. So this names the API path to walk.
            print(f"\n[!] 共 {totals['table'] + totals['group']} 个元素本脚本读不到"
                  f"（表格 ×{totals['table']} · 组合 ×{totals['group']}），"
                  f"上面已按页用 [unread] 标出。")
            print("    它们的文字不在以上输出里。要读，用 python-pptx 直接走："
                  "表格 shape.table.rows[].cells[].text_frame，"
                  "组合 shape.shapes[]（可嵌套）。")
    return 0

if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
