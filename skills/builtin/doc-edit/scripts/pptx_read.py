#!/usr/bin/env python3
"""Read a .pptx: dump slide outline (text of each shape). ultrawork doc-edit skill."""
from __future__ import annotations
import argparse, json, sys

def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def main(argv):
    ap = argparse.ArgumentParser(description="Read slide outline from a .pptx")
    ap.add_argument("file")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args(argv)
    pptx = _require()
    try:
        prs = pptx.Presentation(args.file)
    except Exception as exc:  # noqa: BLE001
        print(f"Error opening {args.file}: {exc}", file=sys.stderr)
        return 1
    slides = []
    for idx, slide in enumerate(prs.slides):
        texts = [s.text for s in slide.shapes if s.has_text_frame and s.text.strip()]
        slides.append({"index": idx, "layout": slide.slide_layout.name, "texts": texts})
    if args.json:
        print(json.dumps({"slides": slides}, ensure_ascii=False, indent=2))
    else:
        for s in slides:
            print(f"=== slide {s['index']} ({s['layout']}) ===")
            for t in s["texts"]:
                print(t)
    return 0

if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
