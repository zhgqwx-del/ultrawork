#!/usr/bin/env python3
"""Edit a .pptx: replace text, add a slide with a title. ultrawork doc-edit skill."""
from __future__ import annotations
import argparse, sys

def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def _replace_in_shape(shape, old, new):
    if not shape.has_text_frame:
        return 0
    n = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                n += 1
    return n

def main(argv):
    ap = argparse.ArgumentParser(description="Edit a .pptx in place (or --out)")
    ap.add_argument("file")
    ap.add_argument("--replace", nargs=2, metavar=("OLD", "NEW"), action="append", default=[])
    ap.add_argument("--add-slide", action="store_true")
    ap.add_argument("--layout", type=int, default=1, help="slide layout index for --add-slide")
    ap.add_argument("--title", default="", help="title text for the new slide")
    ap.add_argument("--out", help="write to this path instead of in place")
    args = ap.parse_args(argv)
    pptx = _require()
    try:
        prs = pptx.Presentation(args.file)
    except Exception as exc:  # noqa: BLE001
        print(f"Error opening {args.file}: {exc}", file=sys.stderr)
        return 1
    total = 0
    for old, new in args.replace:
        for slide in prs.slides:
            for shape in slide.shapes:
                total += _replace_in_shape(shape, old, new)
    added = 0
    if args.add_slide:
        if args.layout < 0 or args.layout >= len(prs.slide_layouts):
            print(f"Bad --layout {args.layout} (have 0..{len(prs.slide_layouts) - 1})", file=sys.stderr)
            return 1
        slide = prs.slides.add_slide(prs.slide_layouts[args.layout])
        if args.title and slide.shapes.title is not None:
            slide.shapes.title.text = args.title
        added = 1
    out = args.out or args.file
    prs.save(out)
    print(f"Saved {out} (replacements: {total}, slides added: {added})")
    return 0

if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
