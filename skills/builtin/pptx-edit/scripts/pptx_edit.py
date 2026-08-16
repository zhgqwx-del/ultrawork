#!/usr/bin/env python3
"""Edit a .pptx: replace text, add a slide with a title. ultrawork pptx-edit skill.

Two measured limits, both stated in SKILL.md 「限制」: replacement is per-run, so a
phrase PowerPoint split across runs does not match; and only top-level shapes with a
text frame are visited, so text inside tables/groups is left untouched.

The second limit is still a limit — this script does not reach into tables or groups
and is not meant to. What changed on 2026-08-16 is that it no longer stays quiet
about it: `replacements: N` still counts only what WAS replaced, and a second line
now counts what it could SEE but not reach. A number that can only go up reads the
same whether it missed nothing or missed nine.

Same day, second change: a run that replaced nothing no longer rewrites the input.
`prs.save()` is a full repackage even when not one character changed.
"""
from __future__ import annotations
import argparse, sys

# ── stdout must be UTF-8 on every platform, and on Windows it is not ─────────
# This script prints whatever text the slides carry, which for this project's users
# is Chinese. Windows encodes a CAPTURED stdout in the machine's ANSI code page
# (cp1252 on a western install), and Python only defaults to UTF-8 from 3.15
# (PEP 686) — CI pins 3.11. Measured: WITHOUT the lines below this script exits 1
# with `UnicodeEncodeError: 'charmap' codec can't encode character '\u7b2c'` the
# moment its output is captured — which is exactly how an agent calls it. The other
# office skills each carry this in a shared module; this one has no shared module,
# so both entry points carry it. Found by the first CI run of this skill's gate,
# not by any local run: it cannot be seen on macOS or Linux.
for _stream in (sys.stdout, sys.stderr):
    if hasattr(_stream, "reconfigure"):
        try:
            _stream.reconfigure(encoding="utf-8")
        except (ValueError, OSError):        # already detached / not reconfigurable
            pass

def _require():
    try:
        import pptx  # noqa: F401
        return pptx
    except ImportError:
        print("Missing dependency: python-pptx (pip install python-pptx)", file=sys.stderr)
        raise SystemExit(1)

def _replace_in_shape(shape, old, new):
    if not shape.has_text_frame:
        return 0
    n = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                n += 1
    return n


def _tf_hits(tf, old):
    n = 0
    try:
        for para in tf.paragraphs:
            for run in para.runs:
                if old in run.text:
                    n += 1
    except Exception:  # noqa: BLE001
        pass
    return n


def _unreachable_hits(shape, old):
    """How many runs match `old` in places --replace cannot reach.

    Same unit as the replaced count above — one per MATCHING RUN, not per
    occurrence — because the two numbers are printed side by side and a reader
    will subtract them.

    Unreachable means: inside a table cell, or anywhere inside a group (at any
    depth, including a table nested in a group). Every property is guarded; this
    runs on files that already proved they can make python-pptx raise.
    """
    n = 0
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    n += _tf_hits(cell.text_frame, old)
            return n
    except Exception:  # noqa: BLE001
        return n
    try:
        subs = getattr(shape, "shapes", None)
        if subs is not None and not shape.has_text_frame:
            for inner in subs:
                try:
                    if inner.has_text_frame:
                        n += _tf_hits(inner.text_frame, old)
                except Exception:  # noqa: BLE001
                    pass
                n += _unreachable_hits(inner, old)
    except Exception:  # noqa: BLE001
        pass
    return n

def main(argv):
    ap = argparse.ArgumentParser(description="Edit a .pptx in place (or --out)")
    ap.add_argument("file")
    ap.add_argument("--replace", nargs=2, metavar=("OLD", "NEW"), action="append", default=[])
    ap.add_argument("--add-slide", action="store_true")
    ap.add_argument("--layout", type=int, default=1, help="slide layout index for --add-slide")
    ap.add_argument("--title", default="", help="title text for the new slide")
    ap.add_argument("--out", help="write to this path instead of in place")
    args = ap.parse_args(argv)
    pptx = _require()
    try:
        prs = pptx.Presentation(args.file)
    except Exception as exc:  # noqa: BLE001
        print(f"Error opening {args.file}: {exc}", file=sys.stderr)
        return 1
    total = 0
    missed: list[tuple[str, int]] = []
    for old, new in args.replace:
        unreachable = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                total += _replace_in_shape(shape, old, new)
                unreachable += _unreachable_hits(shape, old)
        if unreachable:
            missed.append((old, unreachable))
    added = 0
    if args.add_slide:
        # len(prs.slide_layouts) itself reaches through slide_masters[0], which
        # raises IndexError on a .pptx that carries no master — minimal OOXML
        # generators produce those, and the bound check must not be the thing that
        # crashes (measured on a real 10-part .pptx with no ppt/slideMasters).
        try:
            n_layouts = len(prs.slide_layouts)
        except Exception as exc:  # noqa: BLE001
            print(f"Cannot add a slide to {args.file}: it has no slide master ({exc})", file=sys.stderr)
            return 1
        if args.layout < 0 or args.layout >= n_layouts:
            print(f"Bad --layout {args.layout} (have 0..{n_layouts - 1})", file=sys.stderr)
            return 1
        slide = prs.slides.add_slide(prs.slide_layouts[args.layout])
        if args.title and slide.shapes.title is not None:
            slide.shapes.title.text = args.title
        added = 1
    # A run that changed nothing must not rewrite the caller's file. `prs.save()` is
    # a full repackage: measured on the L4 fixture, asking to replace a word the deck
    # does not contain left all 46 parts byte-identical in CONTENT but still produced
    # a different file (sha bfa83db4… → 33223aa3…) — every zip entry reordered and
    # every timestamp reset. That file survived; the point is that nothing was gained
    # by writing it, while a round-trip is exactly where python-pptx drops the parts
    # it does not model. Measured in L4 §四 on the user's real deck: the model ran
    # --replace on a word that lives only in a table, got 0 replacements, and the
    # file was rewritten anyway.
    #
    # An explicit --out is a different request: the caller asked for that artifact to
    # exist, so it is produced even when it is a plain copy.
    out = args.out or args.file
    if args.out or total or added:
        prs.save(out)
        print(f"Saved {out} (replacements: {total}, slides added: {added})")
    else:
        print(f"Unchanged {args.file} (replacements: 0, slides added: 0) — "
              f"没有任何改动，原文件未被重写（重写会重排 zip 条目并重置时间戳）。"
              f"要一份副本用 --out。")
    # `replacements: N` counts only what WAS replaced. On its own it cannot be told
    # apart from "there was nothing else to find", which is the defect SKILL.md
    # 「限制」 warns about in words. State the other half as a number, and name the
    # action that fixes it — a bare status line was dropped 3/3 times in L4 §二,
    # while the one that named a command to run was acted on.
    for old, n in missed:
        print(f"[!] 另有 {n} 处「{old}」在本脚本改不到的位置（表格单元格 / 组合内部），"
              f"未替换。")
    if missed:
        print("    改它们要绕开本脚本：用 python-pptx 走 "
              "shape.table.rows[].cells[].text_frame 与 shape.shapes[]（组合，可嵌套），"
              "并按 run 改（run.text = ...）以保留原格式；"
              "整格赋值（cell.text = ...）会丢掉 run 级格式。")
    return 0

if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
